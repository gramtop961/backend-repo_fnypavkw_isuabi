import os
from io import BytesIO
from typing import List
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class Slide(BaseModel):
    title: str = Field(..., description="Slide title")
    bullets: List[str] = Field(default_factory=list, description="Bullet points for the slide")


class PresentationPayload(BaseModel):
    topic: str = Field("Yoga: History & Advantages")
    slides: List[Slide]
    author: str | None = None


@app.get("/")
def read_root():
    return {"message": "Hello from FastAPI Backend!"}


@app.get("/api/hello")
def hello():
    return {"message": "Hello from the backend API!"}


@app.get("/test")
def test_database():
    """Test endpoint to check if database is available and accessible"""
    response = {
        "backend": "✅ Running",
        "database": "❌ Not Available",
        "database_url": None,
        "database_name": None,
        "connection_status": "Not Connected",
        "collections": []
    }

    try:
        # Try to import database module
        from database import db

        if db is not None:
            response["database"] = "✅ Available"
            response["database_url"] = "✅ Configured"
            response["database_name"] = db.name if hasattr(db, 'name') else "✅ Connected"
            response["connection_status"] = "Connected"

            # Try to list collections to verify connectivity
            try:
                collections = db.list_collection_names()
                response["collections"] = collections[:10]  # Show first 10 collections
                response["database"] = "✅ Connected & Working"
            except Exception as e:
                response["database"] = f"⚠️  Connected but Error: {str(e)[:50]}"
        else:
            response["database"] = "⚠️  Available but not initialized"

    except ImportError:
        response["database"] = "❌ Database module not found (run enable-database first)"
    except Exception as e:
        response["database"] = f"❌ Error: {str(e)[:50]}"

    # Check environment variables
    import os
    response["database_url"] = "✅ Set" if os.getenv("DATABASE_URL") else "❌ Not Set"
    response["database_name"] = "✅ Set" if os.getenv("DATABASE_NAME") else "❌ Not Set"

    return response


@app.post("/generate_pptx")
def generate_pptx(payload: PresentationPayload):
    """Generate a PPTX presentation from provided slides and stream it back."""
    # Lazy import to avoid loading if not used
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(title_layout)
    slide0.shapes.title.text = payload.topic
    subtitle = slide0.placeholders[1]
    subtitle.text = payload.author or "Auto-generated presentation"

    # Content slides
    bullet_layout = prs.slide_layouts[1]
    for s in payload.slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.title
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        first = True
        for b in s.bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(22)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    filename = f"{payload.topic.replace(' ', '_')}.pptx"
    headers = {
        "Content-Disposition": f"attachment; filename=\"{filename}\""
    }

    return StreamingResponse(
        bio,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
